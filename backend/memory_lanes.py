import os
import json
from datetime import datetime
from langchain.memory import ConversationBufferMemory
from langchain.schema import SystemMessage

# Directory structure
BASE_DIR = "chat_histories"
OVERALL_DIR = os.path.join(BASE_DIR, "overall_intelligence")
MEMORY_LANES = {
    "health": os.path.join(OVERALL_DIR, "health_memory"),
    "work": os.path.join(OVERALL_DIR, "work_memory"),
    "journal": os.path.join(OVERALL_DIR, "journal_memory")
}

# Create necessary directories
for directory in [BASE_DIR, OVERALL_DIR] + list(MEMORY_LANES.values()):
    if not os.path.exists(directory):
        os.makedirs(directory)

def process_file_content(file_path):
    """Process and extract text from various file types."""
    import PyPDF2
    import docx
    import pptx
    
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    extracted_text = ""
    
    try:
        if ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    extracted_text += page.extract_text() or ""
        elif ext in ['.docx', '.doc']:
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
        elif ext in ['.pptx', '.ppt']:
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                extracted_text = f.read()
        else:
            return None, f"Unsupported file type: {ext}"
    except Exception as e:
        return None, f"Error processing file: {e}"
    
    if extracted_text.strip():
        return extracted_text, "File processed successfully"
    else:
        return None, "No text extracted from the file"

class MemoryLane:
    def __init__(self, lane_type):
        self.lane_type = lane_type
        self.memory = ConversationBufferMemory()
        self.file_path = os.path.join(MEMORY_LANES[lane_type], f"{lane_type}_memory.json")
        self.load_memory()

    def load_memory(self):
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                history = json.load(f)
                for entry in history:
                    self.memory.save_context(
                        {"input": entry["human"]},
                        {"output": entry["ai"]}
                    )

    def save_memory(self, user_input, ai_response, extracted_text=None):
        history = []
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                history = json.load(f)
        
        entry = {
            "timestamp": datetime.now().isoformat(),
            "human": user_input,
            "ai": ai_response
        }
        if extracted_text:
            entry["document_content"] = extracted_text
        
        history.append(entry)
        with open(self.file_path, 'w', encoding='utf-8') as f:
            json.dump(history, f, ensure_ascii=False, indent=2)
        return True

    def process_file(self, file_path):
        extracted_text, message = process_file_content(file_path)
        if extracted_text:
            self.memory.chat_memory.messages.append(
                SystemMessage(content=f"Document content: {extracted_text}")
            )
            self.save_memory("File uploaded", "Document processed and stored", extracted_text)
            return True, message
        return False, message
    
    def get_memory_history(self):
        """Get the memory history as a list of dictionaries"""
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return []

class OverallIntelligence:
    def __init__(self, client):
        self.client = client
        self.memory_lanes = {
            lane_type: MemoryLane(lane_type)
            for lane_type in MEMORY_LANES.keys()
        }

    def categorize_input(self, user_input):
        try:
            response = self.client.chat.completions.create(
                model="hf:meta-llama/Meta-Llama-3.1-405B-Instruct",
                messages=[
                    {
                        "role": "system",
                        "content": "Categorize the following input into one or more categories: health, work, journal. Respond with only the category names separated by commas."
                    },
                    {"role": "user", "content": user_input}
                ],
                stream=False
            )
            categories = response.choices[0].message.content.strip().lower().split(',')
            return [cat.strip() for cat in categories if cat.strip() in MEMORY_LANES]
        except Exception as e:
            print(f"Error in categorization: {e}")
            return ["journal"]

    def process_file_upload(self, file_path):
        """Process a file upload and store in appropriate memory lanes"""
        categories = self.categorize_input(f"Categorize this file upload: {file_path}")
        results = []
        
        for category in categories:
            success, message = self.memory_lanes[category].process_file(file_path)
            results.append((category, success, message))
            
        return results

    def get_messages_for_query(self, user_input):
        """Get all relevant messages for a query across memory lanes"""
        categories = self.categorize_input(user_input)
        all_relevant_history = []
        
        for category in categories:
            lane = self.memory_lanes[category]
            all_relevant_history.extend(lane.memory.chat_memory.messages)

        messages = [
            {
                "role": "system",
                "content": "You are Tryambakam., accessing multiple memory lanes to provide comprehensive assistance. Consider all relevant historical context when responding."
            }
        ]
        
        for message in all_relevant_history:
            if isinstance(message, SystemMessage):
                messages.append({"role": "system", "content": message.content})
            elif hasattr(message, 'type') and message.type == 'human':
                messages.append({"role": "user", "content": message.content})
            elif hasattr(message, 'type') and message.type == 'ai':
                messages.append({"role": "assistant", "content": message.content})

        messages.append({"role": "user", "content": user_input})
        
        return messages, categories
    
    def save_response(self, user_input, ai_response, categories):
        """Save the response to all relevant memory lanes"""
        for category in categories:
            self.memory_lanes[category].save_memory(user_input, ai_response) 