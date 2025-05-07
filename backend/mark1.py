import openai
import json
import os
from datetime import datetime
import base64
import cv2
from dotenv import load_dotenv
from langchain.chains import ConversationChain
from langchain.llms import OpenAI
from backend.memory_manager import LangChainMemoryManager

# Directory structure
BASE_DIR = "chat_histories"
GENERAL_DIR = os.path.join(BASE_DIR, "general_intelligence")
OVERALL_DIR = os.path.join(BASE_DIR, "overall_intelligence")
MEMORY_LANES = {
    "health": os.path.join(OVERALL_DIR, "health_memory"),
    "work": os.path.join(OVERALL_DIR, "work_memory"),
    "journal": os.path.join(OVERALL_DIR, "journal_memory")
}

CHAT_DIR = os.path.join(GENERAL_DIR, "chat_pages")

# Create necessary directories
for directory in [BASE_DIR, GENERAL_DIR, OVERALL_DIR, CHAT_DIR] + list(MEMORY_LANES.values()):
    if not os.path.exists(directory):
        os.makedirs(directory)

# Initialize the memory manager
memory_manager = LangChainMemoryManager()

class ChatPage:
    def __init__(self, page_id, title):
        self.page_id = page_id
        self.title = title
        self.file_path = os.path.join(CHAT_DIR, f"chat_{page_id}.json")
        self.load_history()
        
    def load_history(self):
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                try:
                    self.history = json.load(f)
                except json.JSONDecodeError:
                    self.history = []
        else:
            self.history = []

    def save_history(self, user_input, ai_response):
        if not self.history:
            self.history = [{
                "timestamp": datetime.now().isoformat(),
                "title": self.title,
                "human": user_input,
                "ai": ai_response
            }]
        else:
            self.history.append({
                "timestamp": datetime.now().isoformat(),
                "human": user_input,
                "ai": ai_response
            })
        
        with open(self.file_path, 'w', encoding='utf-8') as f:
            json.dump(self.history, f, ensure_ascii=False, indent=2)
    
    def process_file(self, file_path):
        """Process a file and add its content to the chat context"""
        content = process_file_content(file_path)
        if content:
            self.save_history(f"[FILE UPLOADED: {os.path.basename(file_path)}]", 
                             f"I've processed the file. You can now ask questions about its content.")
            return True
        return False

class ChatManager:
    def __init__(self):
        self.chat_pages = {}
        self.load_chat_pages()

    def load_chat_pages(self):
        if not os.path.exists(CHAT_DIR):
            return
        
        for file in os.listdir(CHAT_DIR):
            if file.startswith("chat_") and file.endswith(".json"):
                page_id = file[5:-5]  # Remove 'chat_' and '.json'
                try:
                    with open(os.path.join(CHAT_DIR, file), 'r', encoding='utf-8') as f:
                        data = json.load(f)
                        title = data[0]["title"] if data and "title" in data[0] else f"Chat {page_id}"
                        self.chat_pages[page_id] = ChatPage(page_id, title)
                except (json.JSONDecodeError, IndexError, KeyError):
                    # Handle corrupted files
                    self.chat_pages[page_id] = ChatPage(page_id, f"Chat {page_id}")

    def create_chat_page(self, title):
        page_id = str(len(self.chat_pages) + 1)
        chat_page = ChatPage(page_id, title)
        self.chat_pages[page_id] = chat_page
        return page_id

    def get_chat_page(self, page_id):
        return self.chat_pages.get(page_id)

    def list_chat_pages(self):
        return [{"id": page_id, "title": chat_page.title} 
                for page_id, chat_page in self.chat_pages.items()]

    def delete_chat_page(self, page_id):
        """Delete a chat page by ID"""
        if page_id in self.chat_pages:
            # Get the file path before removing from memory
            file_path = self.chat_pages[page_id].file_path
            
            # Remove from memory
            del self.chat_pages[page_id]
            
            # Delete the file if it exists
            if os.path.exists(file_path):
                os.remove(file_path)
            
            return True
        return False

class MemoryLane:
    def __init__(self, lane_id, client):
        self.lane_id = lane_id
        self.file_path = os.path.join(OVERALL_DIR, f"{lane_id}_memory.json")
        self.client = client
        self.load_history()
        
    def load_history(self):
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                try:
                    self.history = json.load(f)
                except json.JSONDecodeError:
                    self.history = []
        else:
            self.history = []
            
    def save_history(self, user_input, ai_response):
        self.history.append({
            "timestamp": datetime.now().isoformat(),
            "human": user_input,
            "ai": ai_response
        })
        
        # Create directory if it doesn't exist
        os.makedirs(os.path.dirname(self.file_path), exist_ok=True)
        
        with open(self.file_path, 'w', encoding='utf-8') as f:
            json.dump(self.history, f, ensure_ascii=False, indent=2)
            
    def get_response(self, user_input):
        # Prepare messages for the API
        system_prompt = f"You are Tryambakam's {self.lane_id} memory lane. You maintain context and information specifically related to {self.lane_id} topics."
        messages = [{"role": "system", "content": system_prompt}]
        
        # Add conversation history (last 10 exchanges)
        recent_history = self.history[-10:] if len(self.history) > 10 else self.history
        for entry in recent_history:
            if "human" in entry and "ai" in entry:
                messages.append({"role": "user", "content": entry["human"]})
                messages.append({"role": "assistant", "content": entry["ai"]})
        
        # Add the current user input
        messages.append({"role": "user", "content": user_input})
        
        try:
            # Get response from API
            from backend.config import client
            completion = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=messages
            )
            
            response = completion.choices[0].message.content
            
            # Save to history
            self.save_history(user_input, response)
            
            return response
        except Exception as e:
            return f"Error: {str(e)}"

class OverallIntelligence:
    def __init__(self, client):
        self.client = client
        self.memory_manager = memory_manager
        
        # Create LangChain conversation chain
        self.chain = ConversationChain(
            llm=OpenAI(temperature=0.7),
            memory=self.memory_manager.overall_memory,
            verbose=True
        )
    
    def process_message(self, message):
        """Process a message and return a response with memory lane categorization"""
        try:
            # Get response from LangChain
            response = self.chain.predict(input=message)
            
            # Categorize the message
            memory_lanes = self.categorize_message(message, response)
            
            # Save to memory lanes
            self.memory_manager.save_to_memory_lanes(message, response, memory_lanes)
            
            return response, memory_lanes
            
        except Exception as e:
            print(f"Error in OverallIntelligence: {e}")
            return "I apologize, but I encountered an error. Please try again later.", ["journal"]
    
    def _save_conversation_history(self):
        # Implement the logic to save conversation history to file
        pass

    def _categorize_message(self, message, ai_response):
        # Implement the logic to categorize the message and determine memory lanes
        pass

# Helper functions
def process_file_content(file_path):
    """Extract content from various file types"""
    if not os.path.exists(file_path):
        return None
        
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    try:
        if ext == '.pdf':
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() or ""
            return text
            
        elif ext in ['.docx', '.doc']:
            import docx
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
            
        elif ext in ['.pptx', '.ppt']:
            import pptx
            prs = pptx.Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
            
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
                
        else:
            return f"Unsupported file type: {ext}"
            
    except Exception as e:
        return f"Error processing file: {str(e)}"

def encode_image_file(image_path):
    """Encode an image file to base64"""
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")

def send_message(self, message):
    # ...existing code...
    
    completion = client.chat.completions.create(
        model="gpt-3.5-turbo",  # Change this from whatever model was here before
        messages=messages,
        max_tokens=1000
    )
    
    # ...rest of the function... 

def categorize_input(user_input):
    """Categorize user input into one or more memory lanes"""
    try:
        response = client.chat.completions.create(
            model="hf:meta-llama/Meta-Llama-3.1-405B-Instruct",
            messages=[
                {
                    "role": "system",
                    "content": "Categorize the following input into one or more categories: health, work, journal. Be generous in your categorization - if it could possibly fit a category, include it. Respond with only the category names separated by commas."
                },
                {"role": "user", "content": user_input}
            ],
            stream=False
        )
        categories = response.choices[0].message.content.strip().lower().split(',')
        valid_categories = [category.strip() for category in categories if category.strip() in ["health", "work", "journal"]]
        
        # Always include at least one category
        if not valid_categories:
            valid_categories = ["journal"]  # Default to journal
            
        print(f"Categorized input as: {valid_categories}")
        return valid_categories
    except Exception as e:
        print(f"Error categorizing input: {e}")
        return ["journal"]  # Default to journal if categorization fails

def save_to_memory_lanes(user_input, ai_response, categories=None):
    """Save conversation to appropriate memory lanes"""
    if not categories:
        categories = categorize_input(user_input)
    
    # If no categories were determined, default to journal
    if not categories:
        categories = ["journal"]
    
    entry = {
        "timestamp": datetime.now().isoformat(),
        "human": user_input,
        "ai": ai_response
    }
    
    # Save to each relevant memory lane
    for category in categories:
        add_to_memory_lane(category, entry)
    
    print(f"Saved conversation to memory lanes: {categories}")
    return categories

def add_to_memory_lane(memory_type, entry):
    """Add an entry to a specific memory lane"""
    # Ensure directory exists
    if not os.path.exists(MEMORY_LANES[memory_type]):
        os.makedirs(MEMORY_LANES[memory_type])
    
    file_path = os.path.join(MEMORY_LANES[memory_type], f"{memory_type}_memory.json")
    try:
        # Try to read existing data
        if os.path.exists(file_path):
            with open(file_path, 'r', encoding='utf-8') as f:
                try:
                    data = json.load(f)
                except json.JSONDecodeError:
                    data = []
        else:
            data = []
        
        # Append new entry
        data.append(entry)
        
        # Write back to file
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
            
        print(f"Successfully saved to {memory_type} memory lane")
    except Exception as e:
        print(f"Error saving to {memory_type} memory lane: {e}") 

# Update ExpertsManager to use LangChain memory
class Expert:
    def __init__(self, expert_id, client):
        # ... existing initialization ...
        self.memory = memory_manager.get_expert_memory(expert_id)
        
        # Create LangChain conversation chain
        self.chain = ConversationChain(
            llm=OpenAI(temperature=0.7),
            memory=self.memory,
            verbose=True
        )
    
    def get_response(self, user_input, include_memory=False):
        # ... existing code ...
        
        # If memory lane access is enabled, include relevant memories
        if include_memory:
            memory_lane = self.info["memory_lane"]
            lane_memory = memory_manager.get_memory_for_lane(memory_lane)
            if lane_memory:
                memory_vars = lane_memory.load_memory_variables({})
                # Add memory context to the prompt
                # ... code to include memory context ...
        
        # Use LangChain for response generation
        response = self.chain.predict(input=user_input)
        
        # Save to history
        self.save_history(user_input, response)
        
        return response 