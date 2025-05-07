import os
import json
from datetime import datetime
from langchain.memory import ConversationBufferMemory
from langchain.schema import HumanMessage, AIMessage, SystemMessage
import PyPDF2
import docx
import pptx

# Directory structure
BASE_DIR = "chat_histories"
GENERAL_DIR = os.path.join(BASE_DIR, "general_intelligence")
CHAT_DIR = os.path.join(GENERAL_DIR, "chat_pages")

# Create necessary directories
for directory in [BASE_DIR, GENERAL_DIR, CHAT_DIR]:
    if not os.path.exists(directory):
        os.makedirs(directory)

class ChatPage:
    def __init__(self, page_id, title):
        self.page_id = page_id
        self.title = title
        self.memory = ConversationBufferMemory()
        self.file_path = os.path.join(CHAT_DIR, f"chat_{page_id}.json")
        self.load_history()

    def load_history(self):
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                history = json.load(f)
                for entry in history:
                    self.memory.save_context(
                        {"input": entry["human"]},
                        {"output": entry["ai"]}
                    )

    def save_history(self, user_input, ai_response):
        history = []
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                history = json.load(f)
        
        # If this is the first entry, include the title
        if not history:
            history.append({
                "timestamp": datetime.now().isoformat(),
                "title": self.title,
                "human": user_input,
                "ai": ai_response
            })
        else:
            history.append({
                "timestamp": datetime.now().isoformat(),
                "human": user_input,
                "ai": ai_response
            })
        
        with open(self.file_path, 'w', encoding='utf-8') as f:
            json.dump(history, f, ensure_ascii=False, indent=2)
    
    def process_file(self, file_path):
        if not os.path.isfile(file_path):
            return False, "File does not exist."
        
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        extracted_text = ""
        
        try:
            if ext == '.pdf':
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        extracted_text += page.extract_text() or ""
            elif ext in ['.docx', '.doc']:
                doc = docx.Document(file_path)
                for para in doc.paragraphs:
                    extracted_text += para.text + "\n"
            elif ext in ['.pptx', '.ppt']:
                prs = pptx.Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            elif ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    extracted_text = f.read()
            else:
                return False, f"Unsupported file type: {ext}"
        except Exception as e:
            return False, f"Error processing file: {e}"
        
        if extracted_text:
            # Store the extracted text as a system message in memory
            self.memory.chat_memory.messages.append(
                SystemMessage(content=f"Document content: {extracted_text}")
            )
            return True, "File uploaded and content added to the chat context."
        else:
            return False, "No text extracted from the file."

    def get_messages_for_api(self):
        """Get messages formatted for the OpenAI API"""
        messages = [{"role": "system", "content": "You are Tryambakam., an AI assistant."}]
        
        # Get the conversation history from memory
        history = []
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                history = json.load(f)
        
        # Add all previous messages from history to the context
        for entry in history:
            if "human" in entry and "ai" in entry:
                messages.append({"role": "user", "content": entry["human"]})
                messages.append({"role": "assistant", "content": entry["ai"]})
        
        # Add any system messages (like uploaded documents)
        for message in self.memory.chat_memory.messages:
            if isinstance(message, SystemMessage):
                messages.append({"role": "system", "content": message.content})
                
        return messages
    
    def get_chat_history(self):
        """Get the chat history as a list of dictionaries"""
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return []

class ChatManager:
    def __init__(self):
        self.chat_pages = {}
        self.load_chat_pages()

    def load_chat_pages(self):
        if not os.path.exists(CHAT_DIR):
            return
        
        for file in os.listdir(CHAT_DIR):
            if file.startswith("chat_") and file.endswith(".json"):
                page_id = file[5:-5]  # Remove 'chat_' and '.json'
                with open(os.path.join(CHAT_DIR, file), 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    title = data[0]["title"] if data and "title" in data[0] else f"Chat {page_id}"
                    self.chat_pages[page_id] = ChatPage(page_id, title)

    def create_chat_page(self, title):
        page_id = str(len(self.chat_pages) + 1)
        # Check if this ID already exists (in case of deletions)
        while page_id in self.chat_pages:
            page_id = str(int(page_id) + 1)
            
        chat_page = ChatPage(page_id, title)
        self.chat_pages[page_id] = chat_page
        return page_id, chat_page

    def get_chat_page(self, page_id):
        return self.chat_pages.get(page_id)

    def list_chat_pages(self):
        return [(page_id, chat_page.title) for page_id, chat_page in self.chat_pages.items()]

    def delete_chat_page(self, page_id):
        """Delete a chat page by ID"""
        if page_id in self.chat_pages:
            # Get the file path
            file_path = self.chat_pages[page_id].file_path
            
            # Remove from memory
            del self.chat_pages[page_id]
            
            # Remove the file if it exists
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                    return True, f"Chat {page_id} deleted successfully"
                except Exception as e:
                    return False, f"Error deleting file: {e}"
            return True, f"Chat {page_id} deleted successfully"
        return False, "Chat not found" 