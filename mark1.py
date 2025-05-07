import openai
import json
import os
from datetime import datetime
from langchain.memory import ConversationBufferMemory
from langchain.schema import HumanMessage, AIMessage, SystemMessage

# New imports for file processing
import PyPDF2
import docx
import pptx

# New imports for vision capabilities
import base64
import cv2
from dotenv import load_dotenv

# Import the intelligence experts system
import intelligence_experts

# Import voice interface
from voice_interface import voice_output

# API configuration
api_key = "glhf_4bcabf37973c831859edc8b224c682f4"
client = openai.OpenAI(
    api_key=api_key,
    base_url="https://glhf.chat/api/openai/v1",
)

# OpenAI client for vision capabilities
vision_api_key = os.getenv("OPENAI_API_KEY", api_key)  # Use environment variable or default
vision_client = openai.OpenAI(api_key=vision_api_key)

# Update directory structure
BASE_DIR = "chat_histories"
GENERAL_DIR = os.path.join(BASE_DIR, "general_intelligence")
OVERALL_DIR = os.path.join(BASE_DIR, "overall_intelligence")
MEMORY_LANES = {
    "health": os.path.join(OVERALL_DIR, "health_memory"),
    "work": os.path.join(OVERALL_DIR, "work_memory"),
    "journal": os.path.join(OVERALL_DIR, "journal_memory")
}

CHAT_DIR = os.path.join(GENERAL_DIR, "chat_pages")

# Create necessary directories
for directory in [BASE_DIR, GENERAL_DIR, OVERALL_DIR, CHAT_DIR] + list(MEMORY_LANES.values()):
    if not os.path.exists(directory):
        os.makedirs(directory)

# Check if there are chat files in the old location and move them to the new location
old_chat_dir = os.path.join(BASE_DIR, "chat_pages")
if os.path.exists(old_chat_dir):
    print(f"Found chat files in old location. Moving to {CHAT_DIR}...")
    import shutil
    for filename in os.listdir(old_chat_dir):
        if filename.startswith("chat_") and filename.endswith(".json"):
            old_path = os.path.join(old_chat_dir, filename)
            new_path = os.path.join(CHAT_DIR, filename)
            shutil.copy2(old_path, new_path)
            print(f"Moved {filename} to new location")
    print("Chat files migration complete")

class ChatPage:
    def __init__(self, page_id, title):
        self.page_id = page_id
        self.title = title
        self.memory = ConversationBufferMemory()
        self.file_path = os.path.join(CHAT_DIR, f"chat_{page_id}.json")
        self.load_history()

    def load_history(self):
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                history = json.load(f)
                for entry in history:
                    self.memory.save_context(
                        {"input": entry["human"]},
                        {"output": entry["ai"]}
                    )

    def save_history(self, user_input, ai_response):
        history = []
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                history = json.load(f)
        
        # If this is the first entry, include the title
        if not history:
            history.append({
                "timestamp": datetime.now().isoformat(),
                "title": self.title,
                "human": user_input,
                "ai": ai_response
            })
        else:
            history.append({
                "timestamp": datetime.now().isoformat(),
                "human": user_input,
                "ai": ai_response
            })
        
        with open(self.file_path, 'w', encoding='utf-8') as f:
            json.dump(history, f, ensure_ascii=False, indent=2)
    
    def process_file(self, file_path):
        if not os.path.isfile(file_path):
            print("File does not exist.")
            return False
        
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        extracted_text = ""
        
        try:
            if ext == '.pdf':
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        extracted_text += page.extract_text() or ""
            elif ext in ['.docx', '.doc']:
                doc = docx.Document(file_path)
                for para in doc.paragraphs:
                    extracted_text += para.text + "\n"
            elif ext in ['.pptx', '.ppt']:
                prs = pptx.Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            elif ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    extracted_text = f.read()
            else:
                print(f"Unsupported file type: {ext}")
                return False
        except Exception as e:
            print(f"Error processing file: {e}")
            return False
        
        if extracted_text:
            # Store the extracted text as a system message in memory with J.A.R.V.I.S. context
            self.memory.chat_memory.messages.append(
                SystemMessage(content=f"Document content: {extracted_text}")
            )
            print("File uploaded and content added to the chat context.")
            return True
        else:
            print("No text extracted from the file.")
            return False

class ChatManager:
    def __init__(self):
        self.chat_pages = {}
        self.load_chat_pages()

    def load_chat_pages(self):
        if not os.path.exists(CHAT_DIR):
            return
        
        for file in os.listdir(CHAT_DIR):
            if file.startswith("chat_") and file.endswith(".json"):
                page_id = file[5:-5]  # Remove 'chat_' and '.json'
                with open(os.path.join(CHAT_DIR, file), 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    title = data[0]["title"] if data and "title" in data[0] else f"Chat {page_id}"
                    self.chat_pages[page_id] = ChatPage(page_id, title)

    def create_chat_page(self, title):
        page_id = str(len(self.chat_pages) + 1)
        chat_page = ChatPage(page_id, title)
        self.chat_pages[page_id] = chat_page
        return page_id

    def get_chat_page(self, page_id):
        return self.chat_pages.get(page_id)

    def list_chat_pages(self):
        if not self.chat_pages:
            print("No chat pages found.")
        else:
            print("\nAvailable Chat Pages:")
            for page_id, chat_page in self.chat_pages.items():
                print(f"Chat ID: {page_id} - Title: {chat_page.title}")

    def delete_chat_page(self, page_id):
        """Delete a chat page by ID"""
        if page_id in self.chat_pages:
            # Get the file path
            file_path = self.chat_pages[page_id].file_path
            
            # Remove from memory
            del self.chat_pages[page_id]
            
            # Remove the file if it exists
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                    return True
                except Exception as e:
                    print(f"Error deleting file: {e}")
                    return False
            return True
        return False

class MemoryLane:
    def __init__(self, lane_type):
        self.lane_type = lane_type
        self.memory = ConversationBufferMemory()
        self.file_path = os.path.join(MEMORY_LANES[lane_type], f"{lane_type}_memory.json")
        self.load_memory()

    def load_memory(self):
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                history = json.load(f)
                for entry in history:
                    self.memory.save_context(
                        {"input": entry["human"]},
                        {"output": entry["ai"]}
                    )

    def save_memory(self, user_input, ai_response, extracted_text=None):
        history = []
        if os.path.exists(self.file_path):
            with open(self.file_path, 'r', encoding='utf-8') as f:
                history = json.load(f)
        
        entry = {
            "timestamp": datetime.now().isoformat(),
            "human": user_input,
            "ai": ai_response
        }
        if extracted_text:
            entry["document_content"] = extracted_text
        
        history.append(entry)
        with open(self.file_path, 'w', encoding='utf-8') as f:
            json.dump(history, f, ensure_ascii=False, indent=2)

    def process_file(self, file_path):
        extracted_text = process_file_content(file_path)
        if extracted_text:
            self.memory.chat_memory.messages.append(
                SystemMessage(content=f"Document content: {extracted_text}")
            )
            self.save_memory("File uploaded", "Document processed and stored", extracted_text)
            return True
        return False

class OverallIntelligence:
    def __init__(self):
        self.memory_lanes = {
            lane_type: MemoryLane(lane_type)
            for lane_type in MEMORY_LANES.keys()
        }

    def categorize_input(self, user_input):
        try:
            response = client.chat.completions.create(
                model="hf:meta-llama/Meta-Llama-3.1-405B-Instruct",
                messages=[
                    {
                        "role": "system",
                        "content": "Categorize the following input into one or more categories: health, work, journal. Respond with only the category names separated by commas."
                    },
                    {"role": "user", "content": user_input}
                ],
                stream=False
            )
            categories = response.choices[0].message.content.strip().lower().split(',')
            return [cat.strip() for cat in categories if cat.strip() in MEMORY_LANES]
        except Exception as e:
            print(f"Error in categorization: {e}")
            return ["journal"]

    def process_input(self, user_input):
        if user_input.lower().startswith('upload '):
            file_path = user_input[7:].strip('"').strip("'")
            categories = self.categorize_input(f"Categorize this file upload: {file_path}")
            for category in categories:
                if self.memory_lanes[category].process_file(file_path):
                    print(f"File processed and stored in {category} memory lane")
            return

        categories = self.categorize_input(user_input)
        all_relevant_history = []
        
        for category in categories:
            lane = self.memory_lanes[category]
            all_relevant_history.extend(lane.memory.chat_memory.messages)

        try:
            messages = [
                {
                    "role": "system",
                    "content": "You are Tryambakam., accessing multiple memory lanes to provide comprehensive assistance. Consider all relevant historical context when responding."
                }
            ]
            
            for message in all_relevant_history:
                if isinstance(message, SystemMessage):
                    messages.append({"role": "system", "content": message.content})
                elif isinstance(message, HumanMessage):
                    messages.append({"role": "user", "content": message.content})
                elif isinstance(message, AIMessage):
                    messages.append({"role": "assistant", "content": message.content})

            messages.append({"role": "user", "content": user_input})

            completion = client.chat.completions.create(
                model="hf:meta-llama/Meta-Llama-3.1-405B-Instruct",
                messages=messages,
                stream=True
            )

            print("\nTryambakam.: ", end='')
            assistant_response = ""
            for chunk in completion:
                if chunk.choices[0].delta.content is not None:
                    content = chunk.choices[0].delta.content
                    assistant_response += content
                    print(content, end='', flush=True)
            print()

            # Add voice output for the complete response
            voice_output.speak(assistant_response)

            for category in categories:
                self.memory_lanes[category].save_memory(user_input, assistant_response)

        except Exception as e:
            print(f"Error getting AI response: {e}")

def process_file_content(file_path):
    """Process and extract text from various file types."""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    extracted_text = ""
    
    try:
        if ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    extracted_text += page.extract_text() or ""
        elif ext in ['.docx', '.doc']:
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
        elif ext in ['.pptx', '.ppt']:
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                extracted_text = f.read()
        else:
            print(f"Unsupported file type: {ext}")
            return None
    except Exception as e:
        print(f"Error processing file: {e}")
        return None
    
    return extracted_text if extracted_text.strip() else None

def print_chat_commands():
    print("\nAvailable commands:")
    print("1. new <title> - Create new chat")
    print("2. switch <chat_id> - Switch to existing chat")
    print("3. upload <file_path> - Upload a file to current chat")
    print("4. list - Show all available chat pages")
    print("5. delete <chat_id> - Delete a chat page")
    print("6. back - Return to main menu")
    print("7. exit - Exit the program")
    print("8. vision - Activate webcam vision analysis")
    print("9. image upload - Analyze an uploaded image")
    print("10. voice - Toggle voice output on/off")

def overall_interface():
    overall_intelligence = OverallIntelligence()
    print("\nOverall Intelligence Activated")
    print("Memory Lanes: Health, Work, Journal")
    print("\nAvailable commands:")
    print("1. upload <file_path> - Upload a file")
    print("2. experts - Access Intelligence Experts")
    print("3. back - Return to main menu")
    print("4. exit - Exit the program")
    print("5. voice - Toggle voice output on/off")
    
    while True:
        user_input = input("\nYou: ").strip()
        if user_input.lower() == 'back':
            break
        if user_input.lower() == 'exit':
            return True
        if user_input.lower() == 'experts':
            # Call the experts system
            if intelligence_experts.start_experts_system(client):
                return True
            print("\nOverall Intelligence Activated")
            print("Memory Lanes: Health, Work, Journal")
            print("\nAvailable commands:")
            print("1. upload <file_path> - Upload a file")
            print("2. experts - Access Intelligence Experts")
            print("3. back - Return to main menu")
            print("4. exit - Exit the program")
            print("5. voice - Toggle voice output on/off")
            continue
        if user_input.lower() == 'voice':
            status = "enabled" if voice_output.toggle_voice() else "disabled"
            print(f"Voice output {status}")
            print("\nOverall Intelligence Activated")
            print("Memory Lanes: Health, Work, Journal")
            print("\nAvailable commands:")
            print("1. upload <file_path> - Upload a file")
            print("2. experts - Access Intelligence Experts")
            print("3. back - Return to main menu")
            print("4. exit - Exit the program")
            print("5. voice - Toggle voice output on/off")
            continue
        overall_intelligence.process_input(user_input)
        
    return False

def chat_interface():
    chat_manager = ChatManager()
    print("\nGeneral Intelligence Activated")
    print_chat_commands()
    
    current_chat = None
    webcam = None
    
    while True:
        if current_chat:
            prompt = f"\nChat {current_chat.page_id} - {current_chat.title}\nYou: "
        else:
            prompt = "\nNo active chat\nCommand: "
            
        user_input = input(prompt).strip()
        
        if user_input.lower() == 'back':
            # Clean up webcam if it was initialized
            if webcam:
                webcam.stop()
                webcam = None
                cv2.destroyAllWindows()
            break
            
        if user_input.lower() == 'exit':
            # Clean up webcam if it was initialized
            if webcam:
                webcam.stop()
                webcam = None
                cv2.destroyAllWindows()
            return True
            
        if user_input.lower() == 'list':
            chat_manager.list_chat_pages()
            print_chat_commands()
            continue
            
        if user_input.lower().startswith('new '):
            title = user_input[4:].strip()
            page_id = chat_manager.create_chat_page(title)
            current_chat = chat_manager.get_chat_page(page_id)
            print(f"Created new chat {page_id}: {title}")
            print_chat_commands()
            
        elif user_input.lower().startswith('switch '):
            page_id = user_input[7:].strip()
            current_chat = chat_manager.get_chat_page(page_id)
            if current_chat:
                print(f"Switched to chat {page_id}: {current_chat.title}")
            else:
                print("Chat not found")
            print_chat_commands()
                
        elif user_input.lower().startswith('upload ') and current_chat:
            file_path = user_input[7:].strip('"').strip("'")
            current_chat.process_file(file_path)
            print_chat_commands()
            
        elif user_input.lower() == 'vision':
            print("\nVision mode activated. Press 'q' or 'Esc' to exit vision mode.")
            try:
                if webcam is None:
                    webcam = WebcamStream()
                
                while True:
                    vision_prompt = input("\nEnter vision prompt (or 'back' to return): ")
                    if vision_prompt.lower() == 'back':
                        break
                        
                    # Read webcam frame and encode it
                    image_base64 = webcam.read_frame(encode=True)
                    if not image_base64:
                        print("Error: Could not read webcam frame.")
                        continue

                    # Generate response
                    print("\nAnalyzing image...")
                    response = generate_vision_response(vision_prompt, image_base64)
                    print(f"\nTryambakam.: {response}")
                    
                    # Display the webcam feed
                    frame = webcam.read_frame()
                    if frame is not None:
                        cv2.imshow("Webcam Feed", frame)
                        if cv2.waitKey(1) in [27, ord("q")]:  # 'Esc' or 'q' to quit
                            break
                
                # Save the last interaction to chat history if in a chat
                if current_chat and response:
                    current_chat.save_history(f"[VISION] {vision_prompt}", response)
                    
            except Exception as e:
                print(f"Error in vision mode: {e}")
            finally:
                print("\nExiting vision mode")
                print_chat_commands()
                
        elif user_input.lower() == 'image upload':
            try:
                image_path = input("\nEnter the path to your image: ")
                if os.path.exists(image_path):
                    image_base64 = encode_image_file(image_path)
                    vision_prompt = input("Enter your prompt about the image: ")
                    
                    print("\nAnalyzing image...")
                    response = generate_vision_response(vision_prompt, image_base64)
                    print(f"\nTryambakam.: {response}")
                    
                    # Save the interaction to chat history if in a chat
                    if current_chat:
                        current_chat.save_history(f"[IMAGE] {vision_prompt}", response)
                else:
                    print("Invalid file path. Please try again.")
            except Exception as e:
                print(f"Error processing image: {e}")
            finally:
                print_chat_commands()
            
        elif user_input.lower() == 'voice':
            status = "enabled" if voice_output.toggle_voice() else "disabled"
            print(f"Voice output {status}")
            print_chat_commands()
            
        elif user_input.lower().startswith('delete '):
            page_id = user_input[7:].strip()
            
            # Check if the chat exists
            chat_to_delete = chat_manager.get_chat_page(page_id)
            if not chat_to_delete:
                print(f"Chat {page_id} not found.")
                print_chat_commands()
                continue
            
            # Ask for confirmation with a more explicit prompt
            print(f"WARNING: You are about to delete Chat {page_id}: '{chat_to_delete.title}'")
            confirm = input("Type 'yes' or 'y' to confirm deletion, anything else to cancel: ").strip().lower()
            
            # Only proceed if user explicitly confirms
            if confirm in ['y', 'yes']:
                # If trying to delete the current chat, clear current_chat first
                if current_chat and current_chat.page_id == page_id:
                    current_chat = None
                
                # Delete the chat page
                if chat_manager.delete_chat_page(page_id):
                    print(f"Chat {page_id} deleted successfully")
                else:
                    print(f"Failed to delete chat {page_id}")
            else:
                print(f"Deletion of Chat {page_id} cancelled")
            
            print_chat_commands()
            
        elif current_chat:
            try:
                # Create a list to hold all messages for the API call
                messages = [{"role": "system", "content": "You are Tryambakam., an AI assistant."}]
                
                # Get the conversation history from memory
                history = []
                if os.path.exists(current_chat.file_path):
                    with open(current_chat.file_path, 'r', encoding='utf-8') as f:
                        history = json.load(f)
                
                # Add all previous messages from history to the context
                for entry in history:
                    if "human" in entry and "ai" in entry:
                        messages.append({"role": "user", "content": entry["human"]})
                        messages.append({"role": "assistant", "content": entry["ai"]})
                
                # Add any system messages (like uploaded documents)
                for message in current_chat.memory.chat_memory.messages:
                    if isinstance(message, SystemMessage):
                        messages.append({"role": "system", "content": message.content})
                
                # Add the current user input
                messages.append({"role": "user", "content": user_input})
                
                completion = client.chat.completions.create(
                    model="hf:meta-llama/Meta-Llama-3.1-405B-Instruct",
                    messages=messages,
                    stream=True
                )
                
                print("\nTryambakam.: ", end='')
                assistant_response = ""
                for chunk in completion:
                    if chunk.choices[0].delta.content is not None:
                        content = chunk.choices[0].delta.content
                        assistant_response += content
                        print(content, end='', flush=True)
                print()
                
                # Add voice output for the complete response
                voice_output.speak(assistant_response)
                
                current_chat.save_history(user_input, assistant_response)
                print_chat_commands()
                
            except Exception as e:
                print(f"Error getting AI response: {e}")
                print_chat_commands()
        else:
            print("No active chat. Please create or switch to a chat first.")
            print_chat_commands()
    
    # Clean up webcam if it was initialized
    if webcam:
        webcam.stop()
        webcam = None
        cv2.destroyAllWindows()
    
    return False

def main_interface():
    while True:
        print("\nWelcome to Tryambakam. Intelligence System")
        print("1. Overall Intelligence (Memory Lanes)")
        print("2. General Intelligence (Chat Pages)")
        print("3. Exit")
        
        choice = input("\nPlease select an option (1-3): ").strip()
        
        if choice == "1":
            if overall_interface():
                break
        elif choice == "2":
            print("\nSwitching to General Intelligence")
            chat_interface()
        elif choice == "3":
            print("\nTryambakam.: Shutting down, sir. Have a good day.")
            break
        else:
            print("Invalid option. Please try again.")

# Webcam handler for vision capabilities
class WebcamStream:
    def __init__(self):
        self.stream = cv2.VideoCapture(0)
        if not self.stream.isOpened():
            raise ValueError("Could not open webcam.")
        self.frame = None

    def read_frame(self, encode=False):
        ret, frame = self.stream.read()
        if not ret:
            return None
        if encode:
            _, buffer = cv2.imencode(".jpeg", frame)
            return base64.b64encode(buffer).decode()
        return frame

    def stop(self):
        self.stream.release()

# Function to encode image from file
def encode_image_file(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")

# Generate response using OpenAI GPT for vision
def generate_vision_response(prompt, image_base64):
    try:
        response = vision_client.chat.completions.create(
            model="gpt-4o",  # Using GPT-4o for vision capabilities
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"},
                        },
                    ],
                }
            ],
            max_tokens=300
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"Error: {e}"

def save_to_memory_lanes(user_input, ai_response, categories=None):
    """Save conversation to appropriate memory lanes"""
    if not categories:
        try:
            # Import the client from config
            from backend.config import client
            
            # Use the Llama model as specified
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",  # Temporary fallback while fixing Llama API access
                messages=[
                    {
                        "role": "system",
                        "content": "Categorize the following input into one or more categories: health, work, journal. Respond with only the category names separated by commas."
                    },
                    {"role": "user", "content": user_input}
                ],
                max_tokens=20
            )
            
            categories_text = response.choices[0].message.content.strip().lower()
            categories = [cat.strip() for cat in categories_text.split(',')]
            print(f"Categories detected: {categories}")
        except Exception as e:
            print(f"Error in categorization: {e}")
            # Default to journal if categorization fails
            categories = ["journal"]
    
    # Save to each category
    for category in categories:
        if category in ["health", "work", "journal"]:
            memory_file = f"chat_histories/overall_intelligence/{category}_memory/{category}_memory.json"
            os.makedirs(os.path.dirname(memory_file), exist_ok=True)
            
            # Load existing memory
            memory_data = []
            if os.path.exists(memory_file):
                try:
                    with open(memory_file, 'r', encoding='utf-8') as f:
                        memory_data = json.load(f)
                except:
                    memory_data = []
            
            # Add new entry
            memory_data.append({
                "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "human": user_input,
                "ai": ai_response
            })
            
            # Save updated memory
            with open(memory_file, 'w', encoding='utf-8') as f:
                json.dump(memory_data, f, indent=2)
    
    return categories

if __name__ == "__main__":
    print("Tryambakam. Dual Intelligence System Initializing...")
    main_interface()